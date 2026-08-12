import os
import re
import urllib.parse
import requests
from bs4 import BeautifulSoup
from marketing_agents import function_tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Helper to perform safe requests with headers
def safe_request(url: str, timeout: int = 5) -> str:
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
    }
    try:
        response = requests.get(url, headers=headers, timeout=timeout)
        if response.status_code == 200:
            return response.text
    except Exception as e:
        pass
    return ""

@function_tool
def search_web_tool(query: str) -> str:
    """
    Search the web for market trends, competitor information, or industry insights.
    
    Args:
        query: The search query to look up.
    """
    encoded_query = urllib.parse.quote_plus(query)
    url = f"https://html.duckduckgo.com/html/?q={encoded_query}"
    html = safe_request(url)
    
    results = []
    if html:
        soup = BeautifulSoup(html, "html.parser")
        links = soup.find_all("a", class_="result__snippet")
        for i, link in enumerate(links[:5]):
            title_node = link.find_previous("a", class_="result__a")
            title = title_node.text.strip() if title_node else "Result"
            snippet = link.text.strip()
            results.append(f"[{i+1}] {title}\nSnippet: {snippet}\n")
            
    if not results:
        # Fallback to high-quality simulated search results based on query keywords
        query_lower = query.lower()
        if "competitor" in query_lower or "vs" in query_lower:
            return (
                f"Simulated search results for: '{query}'\n"
                "[1] Competitor Analysis & Benchmarking Report 2026\n"
                "Snippet: Key competitors in the digital space are focusing on automated workflow solutions, custom CRM integrations, and generative AI content delivery mechanisms. Pricing models shift from flat SaaS towards credit-based consumption.\n"
                "[2] Market Leaders Comparison and Strategic Positioning\n"
                "Snippet: Top 3 competitors show an average customer retention rate of 84%, with marketing spend concentrated heavily on Google search ads (42%), LinkedIn account-based marketing (28%), and content marketing (20%).\n"
                "[3] Competitor Growth Rates & Marketing Channel Distribution\n"
                "Snippet: Competitors are adopting multi-agent customer service and strategy automation. Organic traffic accounts for 35% of acquisitions, indicating strong SEO investment."
            )
        elif "trend" in query_lower or "industry" in query_lower or "market" in query_lower:
            return (
                f"Simulated search results for: '{query}'\n"
                "[1] State of Marketing & Industry Trends 2026\n"
                "Snippet: Interactive content, AI-generated personalization, and localized short-form videos are the primary growth drivers. User privacy changes make zero-party data collection essential for audience building.\n"
                "[2] B2B and B2C Marketing Channel Benchmarks 2026\n"
                "Snippet: Average click-through rates (CTR) on Search ads stand at 3.2%, with conversion rates averaging 2.4%. LinkedIn CPC averages $5.50, while Meta CPC remains around $0.90.\n"
                "[3] Demographic Shifts and Customer Intent Analysis\n"
                "Snippet: Modern buyers demand hyper-personalized outreach. Conversion optimization tools using real-time behavioral insights see a 30% increase in campaign ROI."
            )
        else:
            return (
                f"Simulated search results for: '{query}'\n"
                f"[1] Key insights on {query}\n"
                "Snippet: Recent articles indicate rapid innovation in marketing automation. Major platforms are integrating multi-agent reasoning to improve lead scores and content production speeds.\n"
                "[2] Industry report: Adapting to modern customer acquisition\n"
                "Snippet: Companies using structured content pipelines see a 2x increase in customer touchpoints and highly optimized ad budgets."
            )
            
    return "\n".join(results)

@function_tool
def fetch_url_metadata(url: str) -> str:
    """
    Fetch and analyze the content of a competitor's website or blog post to gather intelligence.
    
    Args:
        url: The absolute URL of the page to scrape.
    """
    if not url.startswith("http"):
        url = "https://" + url
        
    html = safe_request(url)
    if html:
        soup = BeautifulSoup(html, "html.parser")
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
            
        title = soup.title.string.strip() if soup.title else "Untitled Page"
        
        # Get headings
        headings = []
        for h in soup.find_all(["h1", "h2", "h3"])[:10]:
            headings.append(f"{h.name}: {h.text.strip()}")
            
        # Get visible text
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text_content = "\n".join(chunk for chunk in chunks if chunk)[:500]  # Limit characters
        
        return (
            f"Successfully scraped URL: {url}\n"
            f"Page Title: {title}\n\n"
            f"Key Headings:\n" + "\n".join(headings) + "\n\n"
            f"Content Preview:\n{text_content}..."
        )
    else:
        # High quality fallback representing typical landing pages
        domain = urllib.parse.urlparse(url).netloc or url
        return (
            f"Could not connect to live URL (Simulating Scraping for domain: {domain})\n"
            f"Mocked Page Title: {domain.capitalize()} - High Performance Solutions\n\n"
            f"Key Headings:\n"
            f"- H1: Scaling Operations with Next-Generation Workflow Automation\n"
            f"- H2: Why Global Teams Choose Our Platform\n"
            f"- H2: Seamless Integrations with Tools You Already Love\n"
            f"- H3: Secure, Compliant, and Built to Scale\n\n"
            f"Content Preview:\n"
            f"We help enterprises automate complex logic and scale workflows without manual intervention. "
            f"Our software features enterprise-grade security, instant data replication, and real-time dashboard analytics. "
            f"Clients report an average of 40% time savings and 2.5x increase in operational throughput within 30 days of implementation..."
        )

def run_seo_analysis(content: str, keywords: list[str]) -> dict:
    content_lower = content.lower()
    word_count = len(content.split())
    
    analysis = {
        "word_count": word_count,
        "keyword_analysis": {},
        "readability_score": 0,
        "readability_rating": "",
        "recommendations": []
    }
    
    # Calculate density
    for kw in keywords:
        kw_lower = kw.lower()
        count = len(re.findall(r'\b' + re.escape(kw_lower) + r'\b', content_lower))
        density = (count / word_count * 100) if word_count > 0 else 0
        analysis["keyword_analysis"][kw] = {
            "count": count,
            "density_percent": round(density, 2),
            "status": "Good" if 0.5 <= density <= 2.5 else ("Too Low" if density < 0.5 else "Keyword Stuffing")
        }
        
    # Readability estimation (Syllable and sentence count approximations)
    sentences = max(1, len(re.split(r'[.!?]+', content)))
    avg_sentence_len = word_count / sentences
    
    # Estimate complexity based on word length
    long_words = sum(1 for w in content.split() if len(w) > 6)
    long_word_ratio = (long_words / word_count) if word_count > 0 else 0
    
    # Readability index (0-100 scale: higher is easier)
    readability = max(10, min(100, int(100 - (avg_sentence_len * 1.015) - (long_word_ratio * 100 * 0.846))))
    analysis["readability_score"] = readability
    
    if readability > 70:
        analysis["readability_rating"] = "Easy (Suitable for general audience)"
    elif readability > 50:
        analysis["readability_rating"] = "Medium (Suitable for business/tech audience)"
    else:
        analysis["readability_rating"] = "Difficult (Highly technical/academic)"
        
    # Recommendations
    if word_count < 100:
        analysis["recommendations"].append("Content length is very short. Aim for at least 300 words for blog posts/articles.")
    if readability < 45:
        analysis["recommendations"].append("Simplify sentences and use shorter words to improve readability.")
        
    for kw, metrics in analysis["keyword_analysis"].items():
        if metrics["status"] == "Too Low":
            analysis["recommendations"].append(f"Increase frequency of keyword '{kw}' (current density is {metrics['density_percent']}%). Target: 1.0%-2.0%.")
        elif metrics["status"] == "Keyword Stuffing":
            analysis["recommendations"].append(f"Reduce frequency of keyword '{kw}' to avoid SEO penalties (current density is {metrics['density_percent']}%).")
            
    if not analysis["recommendations"]:
        analysis["recommendations"].append("Excellent copy! SEO metrics are fully optimized.")
        
    return analysis

@function_tool
def seo_analyzer_tool(content: str, keywords: list[str]) -> str:
    """
    Analyze marketing copy or content drafts for SEO efficiency, keyword density, and readability.
    Returns a formatted text summary of the audit.
    
    Args:
        content: The text content of the ad, blog post, or page copy.
        keywords: A list of target keywords to evaluate.
    """
    analysis = run_seo_analysis(content, keywords)
    
    # Format as clean markdown text for the LLM
    result_str = [
        "SEO Copy Audit Results:",
        f"- Word Count: {analysis['word_count']} words",
        f"- Readability Score: {analysis['readability_score']}/100 ({analysis['readability_rating']})",
        "Keyword Analysis:"
    ]
    for kw, metrics in analysis['keyword_analysis'].items():
        result_str.append(f"  * '{kw}': Count = {metrics['count']}, Density = {metrics['density_percent']}%, Status = {metrics['status']}")
        
    result_str.append("Recommendations:")
    for rec in analysis['recommendations']:
        result_str.append(f"  - {rec}")
        
    return "\n".join(result_str)

def calculate_budget_forecast(total_budget: float, industry: str, strategy_type: str) -> dict:
    # Industry benchmarks: [Search CPC, Social CPC, Search CTR, Conversions, CPA]
    benchmarks = {
        "saas": {"search_cpc": 4.5, "social_cpc": 1.8, "search_ctr": 0.035, "conv_rate": 0.02, "avg_cpa": 90.0},
        "e-commerce": {"search_cpc": 0.8, "social_cpc": 0.4, "search_ctr": 0.028, "conv_rate": 0.032, "avg_cpa": 15.0},
        "b2b services": {"search_cpc": 3.8, "social_cpc": 2.2, "search_ctr": 0.030, "conv_rate": 0.018, "avg_cpa": 120.0},
        "healthcare": {"search_cpc": 2.5, "social_cpc": 0.9, "search_ctr": 0.032, "conv_rate": 0.024, "avg_cpa": 45.0},
        "real estate": {"search_cpc": 2.0, "social_cpc": 0.7, "search_ctr": 0.025, "conv_rate": 0.015, "avg_cpa": 80.0}
    }
    
    ind_key = industry.lower()
    if ind_key not in benchmarks:
        ind_key = "saas"
    bench = benchmarks[ind_key]
    
    # Adjust factors based on strategy
    splits = {}
    if strategy_type.lower() == "aggressive growth":
        # Shift budget to high-intent search and paid social
        splits = {"Google Search Ads": 0.45, "Paid Social (Meta/LinkedIn)": 0.35, "SEO & Content Marketing": 0.10, "Email/Retargeting": 0.10}
    elif strategy_type.lower() == "organic focus":
        # Shift budget to SEO and content
        splits = {"Google Search Ads": 0.15, "Paid Social (Meta/LinkedIn)": 0.15, "SEO & Content Marketing": 0.55, "Email/Retargeting": 0.15}
    else: # Balanced
        splits = {"Google Search Ads": 0.35, "Paid Social (Meta/LinkedIn)": 0.30, "SEO & Content Marketing": 0.20, "Email/Retargeting": 0.15}
        
    forecasts = {}
    allocated_sum = 0
    total_clicks = 0
    total_conversions = 0
    
    for channel, pct in splits.items():
        amount = total_budget * pct
        allocated_sum += amount
        
        # Calculate impacts
        if channel == "Google Search Ads":
            cpc = bench["search_cpc"]
            clicks = amount / cpc
            convs = clicks * bench["conv_rate"]
        elif channel == "Paid Social (Meta/LinkedIn)":
            cpc = bench["social_cpc"]
            clicks = amount / cpc
            convs = clicks * (bench["conv_rate"] * 0.8) # Social usually has slightly lower conversion
        elif channel == "SEO & Content Marketing":
            # Organic has high upfront cost, long term compounding clicks
            cpc = 1.5 # Blended cost per organic click simulation
            clicks = amount / cpc
            convs = clicks * (bench["conv_rate"] * 1.2) # High intent conversion
        else: # Email/Retargeting
            cpc = 0.3
            clicks = amount / cpc
            convs = clicks * (bench["conv_rate"] * 1.5) # Very high conversion
            
        total_clicks += clicks
        total_conversions += convs
        
        forecasts[channel] = {
            "budget_usd": round(amount, 2),
            "estimated_clicks": int(clicks),
            "estimated_conversions": int(convs),
            "estimated_cpa": round(amount / convs, 2) if convs > 0 else 0
        }
        
    blended_cpa = total_budget / total_conversions if total_conversions > 0 else 0
    # Customer Value Estimator
    avg_ltv = bench["avg_cpa"] * 5
    estimated_revenue = total_conversions * (avg_ltv * 0.3)
    estimated_roi = ((estimated_revenue - total_budget) / total_budget * 100) if total_budget > 0 else 0
    
    return {
        "total_budget": total_budget,
        "industry": industry,
        "strategy": strategy_type,
        "channel_splits": forecasts,
        "aggregate_metrics": {
            "total_clicks": int(total_clicks),
            "total_conversions": int(total_conversions),
            "blended_cpa": round(blended_cpa, 2),
            "estimated_roi_percent": round(estimated_roi, 2),
            "estimated_revenue_usd": round(estimated_revenue, 2)
        }
    }

@function_tool
def budget_forecaster_tool(total_budget: float, industry: str, strategy_type: str) -> str:
    """
    Forecasting tool to distribute ad budgets and calculate key metrics (CTR, CPA, Conversions, ROI) based on benchmarks.
    Returns a formatted text summary of the forecast.
    
    Args:
        total_budget: The total marketing budget (in USD).
        industry: The industry type (SaaS, E-commerce, B2B Services, Healthcare, Real Estate).
        strategy_type: Strategic goal (Aggressive Growth, Balanced, Organic Focus).
    """
    forecast = calculate_budget_forecast(total_budget, industry, strategy_type)
    metrics = forecast["aggregate_metrics"]
    
    # Format as clean markdown text
    result_str = [
        f"Budget Allocation Forecast ({forecast['strategy']} focus for {forecast['industry']}):",
        f"- Total Budget: ${forecast['total_budget']:,.2f}",
        f"- Total Estimated Clicks: {metrics['total_clicks']:,}",
        f"- Total Estimated Conversions: {metrics['total_conversions']:,}",
        f"- Blended CPA: ${metrics['blended_cpa']:.2f}",
        f"- Estimated ROI: {metrics['estimated_roi_percent']}%",
        f"- Expected Revenue: ${metrics['estimated_revenue_usd']:,.2f}",
        "Channel Allocations:"
    ]
    for channel, c_data in forecast['channel_splits'].items():
        result_str.append(
            f"  * {channel}: Budget = ${c_data['budget_usd']:,.2f}, Clicks = {c_data['estimated_clicks']:,}, "
            f"Conversions = {c_data['estimated_conversions']:,}, CPA = ${c_data['estimated_cpa']:.2f}"
        )
        
    return "\n".join(result_str)

def generate_ppt_presentation(
    brand_name: str, 
    industry: str, 
    research_data: str, 
    competitor_data: str, 
    campaign_data: str, 
    content_data: str, 
    budget_data: str
) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)
    
    # Custom Color Palette (Premium Blue/Teal Theme)
    c_primary = RGBColor(12, 35, 64)       # Deep Blue
    c_secondary = RGBColor(0, 150, 136)     # Teal
    c_accent = RGBColor(255, 112, 67)       # Coral Accent
    c_light = RGBColor(245, 246, 248)       # Light Grey Background
    c_dark = RGBColor(33, 33, 33)           # Dark Charcoal
    c_white = RGBColor(255, 255, 255)
    
    # Helper to add backgrounds and custom titles
    def style_slide_title(slide, text, subtitle_text=None, is_dark=False):
        # Add background shape if it's light theme
        if not is_dark:
            background = slide.shapes.add_shape(
                1,  # rectangle
                Inches(0), Inches(0), Inches(13.33), Inches(7.5)
            )
            background.fill.solid()
            background.fill.fore_color.rgb = c_light
            background.line.color.rgb = c_light
            
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = c_white if is_dark else c_primary
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Arial"
            p2.font.size = Pt(16)
            p2.font.color.rgb = c_secondary
            
    # Slide 1: Title Slide (Dark Theme)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_primary
    bg.line.color.rgb = c_primary
    
    # Title box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Strategic Marketing Playbook"
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = c_white
    
    p2 = tf.add_paragraph()
    p2.text = f"Prepared for: {brand_name} ({industry})"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = c_secondary
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Powered by OpenAI Agents Multi-Agent Marketing Team"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(180, 180, 180)
    p3.space_before = Pt(40)
    
    # Helper to add standard text box
    def add_content_slide(title, subtitle, content_str, points_count=4):
        s = prs.slides.add_slide(blank_layout)
        style_slide_title(s, title, subtitle)
        
        tb = s.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Clean clean paragraph builder
        lines = [l.strip() for l in content_str.split("\n") if l.strip()][:points_count * 2]
        for idx, line in enumerate(lines):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = line
            p.font.name = "Arial"
            if line.startswith("-") or line.startswith("*") or line.startswith("•"):
                p.font.size = Pt(16)
                p.font.color.rgb = c_dark
                p.space_after = Pt(12)
                p.level = 1
            else:
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = c_primary
                p.space_before = Pt(14)
                p.level = 0
                
    # Slide 2: Business & Problem Context
    add_content_slide(
        "1. Problem Analysis & Objectives",
        "Defining market obstacles and success parameters",
        f"Business Context:\n- Establishing a scalable brand identity in {industry}\n- Bridging customer trust and high friction conversion paths\nStrategic Objectives:\n- Generate brand affinity, lead flow, and scalable web impressions\n- Decrease Cost-per-Acquisition (CPA) and improve multi-channel customer retention\n- Align product messaging directly to modern buyer behaviors"
    )
    
    # Slide 3: Market Research
    add_content_slide(
        "2. Market Research Insights",
        "Understanding trends, user demographics, and industry shifts",
        research_data if research_data else "Industry Analysis:\n- Growth in automated buyer experiences and custom CRM setups\n- Target demographics emphasize simplicity, clear ROI, and responsive support\n- Organic channels represent high-compounding traffic channels"
    )
    
    # Slide 4: Competitor Intelligence
    add_content_slide(
        "3. Competitor Intelligence",
        "Competitor analysis, pricing structures, and messaging channels",
        competitor_data if competitor_data else "Competitor Landscape:\n- Key competitors offer usage-based pricing or bundle discounts\n- Focus channels are Google search ads and LinkedIn Account-Based Marketing\n- Differentiator needed: High personalization and custom API/integrations"
    )
    
    # Slide 5: Proposed Campaign Plan
    add_content_slide(
        "4. Marketing Strategy & Campaigns",
        "Structured marketing blueprint for target campaigns",
        campaign_data if campaign_data else "Proposed Campaigns:\n- Launch Campaign: Targeting product pain points via high-intent Search Ads\n- Growth Campaign: Direct value delivery via educational whitepapers & email flows\n- SEO Roadmap: Target high-difficulty search terms through topical indexing"
    )
    
    # Slide 6: Multi-Channel Content Strategy
    add_content_slide(
        "5. Content Strategy & SEO",
        "Generating optimized copy blueprints and posting schedules",
        content_data if content_data else "Content Focus:\n- Landing Page Copy: SEO-optimized value statements with CTA integrations\n- Paid Social Ads: Carousel copy focusing on ROI and onboarding convenience\n- Newsletter/Email: Drip campaign targeting product features and reviews"
    )
    
    # Slide 7: Budget Allocations
    add_content_slide(
        "6. Budget Allocation Forecasts",
        "Channel distribution modeling and CPC/CPA projections",
        budget_data if budget_data else "Financial Forecast:\n- Search Channels: Distribute 35% to paid Search (Google Ads)\n- Social Channels: Distribute 30% to Paid Social (LinkedIn/Meta)\n- Compounding Channels: Distribute 20% to SEO and content infrastructure\n- CRM/Email: Distribute 15% to conversion optimization and email retention"
    )
    
    # Slide 8: Multi-Agent Platform Flow (Dark Theme)
    slide_arch = prs.slides.add_slide(blank_layout)
    bg_arch = slide_arch.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg_arch.fill.solid()
    bg_arch.fill.fore_color.rgb = c_primary
    bg_arch.line.color.rgb = c_primary
    
    style_slide_title(slide_arch, "7. Multi-Agent Engine Architecture", "Behind the scenes: OpenAI Agents SDK orchestration", is_dark=True)
    
    # Description text
    tb_arch = slide_arch.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
    tf_arch = tb_arch.text_frame
    tf_arch.word_wrap = True
    
    p = tf_arch.paragraphs[0]
    p.text = "System Orchestration Details:"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = c_secondary
    p.space_after = Pt(14)
    
    architectures = [
        "1. Triage Agent (Marketing Director) - Evaluates prompt, coordinates handoffs and returns final response",
        "2. Market Research Agent - Executes duckduckgo search and extracts industry trend benchmarks",
        "3. Competitor Analysis Agent - Scrapes competitor urls and maps competitive positioning matrices",
        "4. Campaign Planner Agent - Constructs high-intent campaign funnels and timeline structures",
        "5. Content Strategist Agent - Leverages SEO tools to optimize keyword density, hooks and titles",
        "6. Optimisation Advisor Agent - Analyzes metrics, runs budget models and delivers final performance optimizations"
    ]
    
    for arc in architectures:
        p_arc = tf_arch.add_paragraph()
        p_arc.text = arc
        p_arc.font.name = "Arial"
        p_arc.font.size = Pt(16)
        p_arc.font.color.rgb = c_white
        p_arc.space_after = Pt(10)
        p_arc.level = 1
        
    # Slide 9: Conclusion / Summary
    add_content_slide(
        "8. Playbook Highlights & Actions",
        "Immediate steps to launch the marketing campaigns",
        "Next Steps:\n- Set up Google Ads targeting high-intent industry-related keywords\n- Create content calendars based on Content Strategist keyword analysis\n- Deploy automated budget trackers modeled by the Budget Forecaster Tool\n- Schedule weekly multi-agent audits to review performance conversion goals"
    )
    
    # Save presentation
    dir_path = "C:\\Users\\Tirth\\.gemini\\antigravity\\scratch\\ai-marketing-strategy-manager"
    if not os.path.exists(dir_path):
        os.makedirs(dir_path)
    file_path = os.path.join(dir_path, "marketing_strategy_presentation.pptx")
    prs.save(file_path)
    
    return file_path

@function_tool
def ppt_presentation_tool(
    brand_name: str, 
    industry: str, 
    research_data: str, 
    competitor_data: str, 
    campaign_data: str, 
    content_data: str, 
    budget_data: str
) -> str:
    """
    Generate a professional 10-12 slide PowerPoint presentation summarizing the marketing strategy.
    Saves the file to the workspace directory.
    
    Args:
        brand_name: Name of the brand.
        industry: Industry sector.
        research_data: Research summaries.
        competitor_data: Competitor analysis.
        campaign_data: Campaign plans.
        content_data: Content strategies.
        budget_data: Budget metrics.
    """
    return generate_ppt_presentation(brand_name, industry, research_data, competitor_data, campaign_data, content_data, budget_data)

def generate_marketing_report(brand_name: str, analysis_results: str) -> str:
    dir_path = "C:\\Users\\Tirth\\.gemini\\antigravity\\scratch\\ai-marketing-strategy-manager"
    if not os.path.exists(dir_path):
        os.makedirs(dir_path)
    file_path = os.path.join(dir_path, "marketing_strategy_report.md")
    
    report_content = f"""# Marketing Strategy Report for {brand_name}

## 1. Executive Summary & Problem Analysis
This report was compiled by the **AI Marketing Strategy Manager** multi-agent pipeline. It documents industry trends, target demographics, competitor benchmarks, content blueprints, and financial forecast models.

## 2. Multi-Agent Playbook Details
{analysis_results}

## 3. Optimization Recommendations
- Focus media spend on high-intent search ads to drive immediate conversions.
- Continuously optimize organic content around key SEO keywords to build compounding value.
- Leverage weekly agentic audits to monitor acquisition metrics and shift budgets dynamically.
"""
    
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(report_content)
        
    return file_path

@function_tool
def marketing_report_tool(brand_name: str, analysis_results: str) -> str:
    """
    Compile analysis details and generate a formatted markdown marketing report file.
    Saves the file to the workspace directory.
    
    Args:
        brand_name: Name of the brand.
        analysis_results: The text content summarizing all the marketing strategy details.
    """
    return generate_marketing_report(brand_name, analysis_results)
